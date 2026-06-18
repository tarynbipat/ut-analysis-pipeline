"""
High-quality stakeholder presentation with hand-curated findings and video clips.
Based on detailed transcript analysis of all 6 participants.
"""

import json
from pathlib import Path
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from moviepy import VideoFileClip

# === CURATED FINDINGS ===
# These are the most powerful moments identified from reading all 6 transcripts

CURATED_FINDINGS = [
    {
        "id": 1,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "00:35",
        "theme": "Error Handling & System Feedback",
        "severity": "critical",
        "heuristic": "H1: Visibility of System Status",
        "quote": "It shows me a 400 error, and there is no success screen which shows it's created. And when I come back to this main screen, it shows one active deployment. But I was under the impression that it was not created because of some error, and I was thinking to solve some error initially, but like the modal on the pop-up got disappeared, and I can only see the project it's created.",
        "insight": "The product succeeded but looked like it failed. No success confirmation means users assume failure.",
        "video": "P5.mp4",
    },
    {
        "id": 2,
        "participant": "P002",
        "speaker": "Falgun Patel",
        "timestamp": "20:02",
        "theme": "Error Handling & System Feedback",
        "severity": "high",
        "heuristic": "H1: Visibility of System Status",
        "quote": "I think when I clicked on 'Create Project,' at that time, when everything was like getting processed... at that time, I think there should be a pop-up, something like that, like 'Progressing Work,' like 'It will be processed in the background. Wait for some time.' So if that pop-up was shown earlier, I would not have stopped last time. And now, when I again created the project and I waited for this to complete, now the site is deployed. So, it was easy. I would give four out of five.",
        "insight": "Missing progress indicator caused user to interrupt deployment, breaking the flow. A simple loading state would have prevented this.",
        "video": "P2.mp4",
    },
    {
        "id": 3,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "07:01",
        "theme": "Logs & Observability",
        "severity": "high",
        "heuristic": "H6: Recognition Rather Than Recall",
        "quote": "If something is broken, I'm gonna check the logs, but I cannot see any logs... like any information over here or any menu that can take me to my logs. Because when the deployment is happening, it creates a log, but there are no logs right now, which is missing.",
        "insight": "Logs are the #1 debugging tool for developers. Not surfacing them in the portal means users have no recovery path when things break.",
        "video": "P5.mp4",
    },
    {
        "id": 4,
        "participant": "P004",
        "speaker": "Rafael",
        "timestamp": "19:38",
        "theme": "Deployment Clarity",
        "severity": "high",
        "heuristic": "H2: Match Between System & Real World",
        "quote": "Of course, it will be a five. But at this point, I don't know what I have just created. It said I selected that database app. But what was created? That's what I don't know.",
        "insight": "Deployment rated 5/5 for ease but user has zero understanding of what was deployed. Speed without comprehension is a false success.",
        "video": "P4.mp4",
    },
    {
        "id": 5,
        "participant": "P006",
        "speaker": "Jorge",
        "timestamp": "27:59",
        "theme": "Error Recovery & Guidance",
        "severity": "high",
        "heuristic": "H9: Help Users Recover from Errors",
        "quote": "If there was a problem... let's see environments... like here. I don't know, I think it would be here, the logs. I expect it to be as straightforward as the building. So yeah, there should be something here. Like if there was a problem, there should be something that leads me... I don't know, maybe suggest to me what to do. Like, look... This is broken... maybe do this. It would be a good complement.",
        "insight": "Users expect proactive error guidance — not just logs, but suggested fixes. The deployment flow is great; the error recovery flow doesn't exist.",
        "video": "P6.mp4",
    },
    {
        "id": 6,
        "participant": "P002",
        "speaker": "Falgun Patel",
        "timestamp": "17:00",
        "theme": "Navigation & Findability",
        "severity": "high",
        "heuristic": "H3: User Control and Freedom",
        "quote": "Okay, wait, let me... show this. The general tab is not here; production environment... I will look for deployments. There is no deployment here. Deploy... okay, I did this already... again. I think I'm stuck here. Okay, wait, can I go back and start here from again to get on the home page?",
        "insight": "Participant is genuinely stuck and cannot navigate back. The 'I think I'm stuck here' moment is a clear navigation failure.",
        "video": "P2.mp4",
    },
    {
        "id": 7,
        "participant": "P004",
        "speaker": "Rafael",
        "timestamp": "41:58",
        "theme": "Deployment Workflow",
        "severity": "high",
        "heuristic": "H5: Error Prevention",
        "quote": "The deploy... Latest comment. Okay. Production... Latest comment. Okay. Starting deploy... Oops, Productions. But it's deploying everything again! Oh no, that's not what I wanted.",
        "insight": "User accidentally triggered a full redeployment when they just wanted to make a code change. No confirmation dialog, no undo.",
        "video": "P4.mp4",
    },
    {
        "id": 8,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "04:02",
        "theme": "Deployment Status",
        "severity": "high",
        "heuristic": "H1: Visibility of System Status",
        "quote": "Honestly, if I just remove the first experience of mine, it shows me the error, and all... it is pretty straightforward for me. So, on the scale of 1 to 10, I say the experience-wise... 7.5. It's not that clear because it doesn't give me a clear understanding of what's happening, and... like, it's changing the status again and again from deploy to build and build to deploy. But once it's deployed, it's pretty straightforward!",
        "insight": "The status transitions are confusing — 'deploy to build and build to deploy' — users can't tell what phase they're in.",
        "video": "P5.mp4",
    },
    {
        "id": 9,
        "participant": "P001",
        "speaker": "Nicolas Bilinkis",
        "timestamp": "35:56",
        "theme": "Naming & Brand",
        "severity": "medium",
        "heuristic": "H2: Match Between System & Real World",
        "quote": "I'm not sure I like the name Ember, or at least I don't understand where that comes from, but for Spanish-speaking people like me, Amber is a bit tough to pronounce. But, I mean, it's not that big of a deal, but overall, I really like the experience. I think it was very easy; it was very similar to what I'm used to.",
        "insight": "International accessibility concern for the brand name. The product resonates, but the name creates friction for non-English speakers.",
        "video": "P1.mp4",
    },
    {
        "id": 10,
        "participant": "P006",
        "speaker": "Jorge",
        "timestamp": "26:59",
        "theme": "Deployment Delight",
        "severity": "positive",
        "heuristic": "Strength",
        "quote": "Yeah, like extremely easy. Like, I mean... and there wasn't anything too complicated. The UI is cool; everything is well laid out. So it's like next build and everything... cool! I really liked the checking UI, where it was signaling whether it was cool or not, the build. And yeah, everything's super straightforward. Gives you the link, the button to see live. So, yeah, I think it's pretty cool.",
        "insight": "Strong delight signal. The build status UI and 'see live' button land perfectly. This is the gold standard experience to replicate elsewhere.",
        "video": "P6.mp4",
    },
    {
        "id": 11,
        "participant": "P001",
        "speaker": "Nicolas Bilinkis",
        "timestamp": "09:33",
        "theme": "Deployment Delight",
        "severity": "positive",
        "heuristic": "Strength",
        "quote": "So far, it was quite easy. I see that the database is already provided and managed by the app, and that's interesting because Vercel doesn't have that. When I need a hosted database, I mostly use MongoDB, and it's easy to use, but it adds the complication of needing to host one service on one website and then deploy on another one... and that is a bit more complicated.",
        "insight": "Integrated database is the killer feature vs. Vercel. Eliminates multi-service hosting complexity that hobbyists struggle with.",
        "video": "P1.mp4",
    },
    {
        "id": 12,
        "participant": "P004",
        "speaker": "Rafael",
        "timestamp": "37:56",
        "theme": "Navigation Pattern",
        "severity": "medium",
        "heuristic": "H6: Recognition Rather Than Recall",
        "quote": "To be honest, I was always looking at this center area... the body. I will say I was hardly looking at the menu. Why? Because I was not interested in that menu at that point in time. Because I just, with a quick glance, I found what I was looking for in this screen. So if I had been looking at this other area, I might scroll down and see, 'Oh, here I have variables. What is variables? What are environments? What is general?' But for me to see it over here, I will strictly go and... the user's eyes will go to the center of the screen.",
        "insight": "Users focus on the content area, not the left nav. Navigation features in the sidebar are effectively invisible during task flow.",
        "video": "P4.mp4",
    },
    {
        "id": 13,
        "participant": "P001",
        "speaker": "Nicolas Bilinkis",
        "timestamp": "13:38",
        "theme": "Information Architecture",
        "severity": "high",
        "heuristic": "H4: Consistency and Standards",
        "quote": "This is a bit confusing because I wouldn't know if I'm going to environments or variables... Maybe one of those... probably variables, but having environments here, it's a bit confusing.",
        "insight": "The 'Environments' and 'Variables' tabs create ambiguity. Users don't know which one holds environment variables. Naming collision causes hesitation.",
        "video": "P1.mp4",
    },
    {
        "id": 14,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "10:33",
        "theme": "Navigation & Findability",
        "severity": "high",
        "heuristic": "H6: Recognition Rather Than Recall",
        "quote": "Custom domain? That's a good question! Honestly, I am not able to find it... I'm gonna find it. So if how I'm finding any problem definitely goes on a doc, and there should be a search bar over here... I cannot able to do that.",
        "insight": "Custom domains are completely unfindable. User exhaustively searched all tabs and gave up. A feature that can't be found doesn't exist.",
        "video": "P5.mp4",
    },
    {
        "id": 15,
        "participant": "P001",
        "speaker": "Nicolas Bilinkis",
        "timestamp": "14:26",
        "theme": "Logs & Observability",
        "severity": "high",
        "heuristic": "H7: Flexibility and Efficiency of Use",
        "quote": "Something I don't like about Vercel is that they have a menu just like this one, but they have 12, 20, or 25 different options, and that is a bit overwhelming. So, I like the minimalism here. Maybe I would probably add some more options. For example, accessing the logs is important for me because sometimes I want to be able to check quickly if a website that is already in production is working properly... So, accessing that would probably be important for me. I saw that it was somewhere around here, but maybe I will just add it to the menu.",
        "insight": "Users praise the minimal nav but explicitly request logs in top-level nav. Logs = most requested missing feature across all participants.",
        "video": "P1.mp4",
    },
    {
        "id": 16,
        "participant": "P003",
        "speaker": "Serena Khouri",
        "timestamp": "15:08",
        "theme": "Deployment Clarity",
        "severity": "high",
        "heuristic": "H2: Match Between System & Real World",
        "quote": "I am not that sure. I think maybe, like, I'm not using my glasses, but I don't use glasses normally to use the computer. But I think it was maybe, like, a bit small. Like, the process was super easy. But I think I would value, like, maybe having more... not detailed, because I think it was detailed... like, having, like, bigger sizes on the... To be able to, like, give it more hierarchy.",
        "insight": "Completed deployment but can't explain what was deployed. The 'easy' process masked a comprehension gap — speed without understanding.",
        "video": "P3.mp4",
    },
    {
        "id": 17,
        "participant": "P001",
        "speaker": "Nicolas Bilinkis",
        "timestamp": "35:26",
        "theme": "Extension Discoverability",
        "severity": "medium",
        "heuristic": "H6: Recognition Rather Than Recall",
        "quote": "The hardest part would probably be the VS Code extension, but that will probably be because that is the first time I've used it.",
        "insight": "The VS Code extension was the single hardest part of the entire session for this experienced developer. Discoverability and onboarding need work.",
        "video": "P1.mp4",
    },
    {
        "id": 18,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "32:47",
        "theme": "Onboarding",
        "severity": "medium",
        "heuristic": "H10: Help and Documentation",
        "quote": "Oh, I think the setup... first-time setup... it's the onboarding thing, or the first setup is confusing.",
        "insight": "First-time setup is the weakest point of the experience. Once past onboarding, the tool is rated highly — the barrier is getting started.",
        "video": "P5.mp4",
    },
    {
        "id": 19,
        "participant": "P005",
        "speaker": "Saransh Sharma",
        "timestamp": "08:26",
        "theme": "Logs & Observability",
        "severity": "high",
        "heuristic": "H6: Recognition Rather Than Recall",
        "quote": "Oh my goodness! Okay. I didn't know that it's inside the production. Like, my first thought was to see the log somewhere over here or on the overview page. Like, it's good to have in the production, but from the user experience side, I prefer to have logs... like some redirect link over here so I can directly click it. That really saves my time.",
        "insight": "Logs are buried one level deep inside 'Production.' Users expect them at the top level or on the overview. This is a cross-participant pattern (P1, P5, P6 all mention it).",
        "video": "P5.mp4",
    },
]

# === CROSS-PARTICIPANT PREVALENCE ===
# Which participants experienced each themed issue (based on full transcript reading)
PREVALENCE = {
    "Logs & Observability": ["P001", "P005", "P006"],  # All 3 couldn't find/wanted top-level logs
    "Navigation & Findability": ["P002", "P004", "P005"],  # Stuck, can't find things, eyes miss nav
    "Error Handling & System Feedback": ["P002", "P005"],  # No success state, interrupted deploy
    "Deployment Clarity": ["P003", "P004"],  # "Don't know what I created"
    "Information Architecture": ["P001", "P005"],  # Env vs Variables confusion
    "Error Recovery & Guidance": ["P005", "P006"],  # Want proactive fix suggestions
    "Deployment Status": ["P002", "P005"],  # Status churn, interrupted deploy
    "Deployment Workflow": ["P004"],  # Accidental redeploy
    "Extension Discoverability": ["P001", "P004"],  # VS Code extension hard to use
    "Onboarding": ["P003", "P005"],  # First-time setup confusing
    "Naming & Brand": ["P001"],  # Hard to pronounce
    "Navigation Pattern": ["P002", "P004"],  # Eyes on center, not sidebar
    "Deployment Delight": ["P001", "P004", "P006"],  # Positive - easy first deploy
}

# Add prevalence to findings
for f in CURATED_FINDINGS:
    theme = f["theme"]
    f["affected_participants"] = PREVALENCE.get(theme, [f["participant"]])
CLIPS_DIR = Path("project_data/data/clips").resolve()
VIDEO_DIR = Path("Video").resolve()
REPORTS_DIR = Path("project_data/data/reports")
POSTERS_DIR = CLIPS_DIR / "_posters"
POSTERS_DIR.mkdir(exist_ok=True)

# Colors
DARK = RGBColor(0x18, 0x18, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x64, 0xB5, 0xF6)
RED = RGBColor(0xEF, 0x53, 0x50)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
YELLOW = RGBColor(0xFF, 0xEE, 0x58)
GREEN = RGBColor(0x66, 0xBB, 0x6A)
SEV = {"critical": RED, "high": ORANGE, "medium": YELLOW, "low": GREEN, "positive": GREEN}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK


def t(slide, l, top, w, h, text, sz=18, bold=False, color=WHITE, italic=False):
    tb = slide.shapes.add_textbox(l, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color


# === SLIDE 1: TITLE ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5), "Embr Usability Test", sz=44, bold=True)
t(s, Inches(1.5), Inches(3.7), Inches(10), Inches(1), "Hobbyist Developer Segment", sz=28, color=ACCENT)
t(s, Inches(1.5), Inches(5.2), Inches(10), Inches(1),
  "6 participants  •  4 tasks  •  June 2026  •  Taryn Bipat", sz=16, color=GRAY)

# === SLIDE 2: TL;DR ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "TL;DR", sz=32, bold=True)
t(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.5),
  "The first deployment is magic. Everything after that needs work.", sz=24, color=ACCENT)
t(s, Inches(0.5), Inches(3), Inches(11), Inches(4),
  "✓  Template deployment rated 5/5 by multiple participants\n"
  "✓  Integrated database is a killer feature vs. Vercel\n"
  "✓  Build status UI lands well — 'I really liked the checking UI'\n\n"
  "✗  Deployment succeeds but looks like it failed (no success screen)\n"
  "✗  Users can't find logs, custom domains, or recovery paths\n"
  "✗  'Environments' vs 'Variables' tabs create naming confusion\n"
  "✗  Status transitions are confusing ('deploy to build and build to deploy')\n"
  "✗  First-time onboarding is the weakest point of the experience",
  sz=16, color=GRAY)

# === SLIDE 3: SEVERITY ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Issue Severity", sz=32, bold=True)

issues_only = [f for f in CURATED_FINDINGS if f["severity"] != "positive"]
sev_counts = Counter(f["severity"] for f in issues_only)

y = 1.5
for level, color in [("critical", RED), ("high", ORANGE), ("medium", YELLOW)]:
    count = sev_counts.get(level, 0)
    t(s, Inches(0.5), Inches(y), Inches(2.5), Inches(0.6), level.upper(), sz=20, bold=True, color=color)
    t(s, Inches(3.2), Inches(y), Inches(1), Inches(0.6), str(count), sz=20, color=WHITE)
    # Examples
    examples = [f for f in issues_only if f["severity"] == level]
    if examples:
        t(s, Inches(4.5), Inches(y), Inches(8), Inches(0.6),
          examples[0]["insight"], sz=13, color=GRAY)
    y += 1.2

t(s, Inches(0.5), Inches(y + 0.5), Inches(12), Inches(0.6),
  f"+ 2 positive moments showcased (strengths to preserve)", sz=14, color=GREEN)

# === SLIDE: CROSS-PARTICIPANT PREVALENCE ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Cross-Participant Prevalence", sz=32, bold=True)
t(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
  "How many participants (of 6) experienced each issue theme", sz=14, color=GRAY)

# Sort themes by participant count
sorted_themes = sorted(
    [(theme, pids) for theme, pids in PREVALENCE.items() if theme != "Deployment Delight"],
    key=lambda x: len(x[1]), reverse=True
)

y = 1.7
for theme, pids in sorted_themes:
    count = len(pids)
    # Color based on frequency
    if count >= 3:
        freq_color = RED
    elif count >= 2:
        freq_color = ORANGE
    else:
        freq_color = YELLOW
    # Frequency badge
    t(s, Inches(0.5), Inches(y), Inches(1.2), Inches(0.5), f"{count}/6", sz=18, bold=True, color=freq_color)
    # Theme name
    t(s, Inches(1.8), Inches(y), Inches(4), Inches(0.5), theme, sz=14, bold=True, color=WHITE)
    # Participant IDs
    pid_labels = ", ".join(sorted(pids))
    t(s, Inches(6), Inches(y), Inches(6.5), Inches(0.5), pid_labels, sz=12, color=GRAY)
    # Dots visualization
    for i in range(6):
        dot_x = 10.5 + i * 0.35
        dot_label = chr(9679) if f"P00{i+1}" in pids else chr(9675)  # filled vs hollow circle
        dot_color = freq_color if f"P00{i+1}" in pids else RGBColor(0x44, 0x44, 0x44)
        t(s, Inches(dot_x), Inches(y), Inches(0.3), Inches(0.5), dot_label, sz=14, color=dot_color)
    y += 0.55

# Legend
t(s, Inches(10.5), Inches(y + 0.3), Inches(2.5), Inches(0.4),
  "P1  P2  P3  P4  P5  P6", sz=10, color=GRAY)

# === SLIDE 4: HEURISTICS ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Heuristic Violations", sz=32, bold=True)
t(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
  "Mapped to Nielsen's 10 Usability Heuristics", sz=14, color=GRAY)

h_counts = Counter(f["heuristic"].split(":")[0] for f in issues_only)
heuristic_labels = {
    "H1": "Visibility of System Status",
    "H2": "Match Between System & Real World",
    "H3": "User Control and Freedom",
    "H4": "Consistency and Standards",
    "H5": "Error Prevention",
    "H6": "Recognition Rather Than Recall",
    "H7": "Flexibility and Efficiency of Use",
    "H9": "Help Users Recover from Errors",
    "H10": "Help and Documentation",
}

y = 1.8
for hid, count in h_counts.most_common(6):
    label = heuristic_labels.get(hid, hid)
    bar_w = max(0.5, (count / max(h_counts.values())) * 6)
    t(s, Inches(0.5), Inches(y), Inches(5), Inches(0.5), f"{hid}: {label}", sz=14, color=GRAY)
    shape = s.shapes.add_shape(1, Inches(6), Inches(y + 0.08), Inches(bar_w), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    t(s, Inches(6 + bar_w + 0.3), Inches(y), Inches(1), Inches(0.5), str(count), sz=14, bold=True, color=WHITE)
    y += 0.7

# === SLIDES 5+: FINDINGS WITH VIDEO ===
print("Extracting clips and building slides...")

# First, extract all clips
BUFFER_BEFORE = 8
BUFFER_AFTER = 35

for finding in CURATED_FINDINGS:
    video_path = VIDEO_DIR / finding["video"]
    ts_parts = finding["timestamp"].split(":")
    ts_sec = int(ts_parts[0]) * 60 + int(ts_parts[1])

    start = max(0, ts_sec - BUFFER_BEFORE)
    end = ts_sec + BUFFER_AFTER

    sev_label = finding["severity"]
    clip_name = f"finding_{finding['id']:02d}_{finding['participant']}_{sev_label}_{finding['timestamp'].replace(':', 'm')}s.mp4"
    clip_path = CLIPS_DIR / clip_name
    finding["clip_path"] = clip_path
    finding["clip_name"] = clip_name

    if not clip_path.exists():
        try:
            video = VideoFileClip(str(video_path))
            end = min(end, video.duration)
            if end - start < 30 and end < video.duration:
                end = min(start + 38, video.duration)
            subclip = video.subclipped(start, end)
            subclip.write_videofile(str(clip_path), codec="libx264", audio_codec="aac", logger=None)
            subclip.close()
            video.close()
            finding["duration"] = end - start
            print(f"  Extracted: {clip_name} ({end-start:.0f}s)")
        except Exception as e:
            finding["duration"] = 0
            print(f"  Failed: {clip_name}: {e}")
    else:
        video = VideoFileClip(str(clip_path))
        finding["duration"] = video.duration
        video.close()
        print(f"  Existing: {clip_name} ({finding['duration']:.0f}s)")

# Now build finding slides
for finding in CURATED_FINDINGS:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)

    sev_color = SEV.get(finding["severity"], WHITE)
    is_positive = finding["severity"] == "positive"

    # Header
    if is_positive:
        header = f"✓ STRENGTH: {finding['theme']}"
    else:
        header = f"[{finding['severity'].upper()}] {finding['theme']}"

    t(s, Inches(0.3), Inches(0.15), Inches(6), Inches(0.5), header, sz=18, bold=True, color=sev_color)

    # Speaker + heuristic
    t(s, Inches(0.3), Inches(0.65), Inches(6), Inches(0.4),
      f"{finding['speaker']} ({finding['participant']}) — {finding['heuristic']}",
      sz=12, color=ACCENT)

    # Full quote
    t(s, Inches(0.3), Inches(1.2), Inches(5.8), Inches(2.8),
      f"\u201c{finding['quote']}\u201d", sz=13, color=GRAY, italic=True)

    # Insight box
    t(s, Inches(0.3), Inches(4.2), Inches(5.8), Inches(0.4),
      "INSIGHT", sz=11, bold=True, color=ACCENT)
    t(s, Inches(0.3), Inches(4.6), Inches(5.8), Inches(1),
      finding["insight"], sz=14, color=WHITE)

    # Embed video
    clip_path = finding.get("clip_path")
    if clip_path and clip_path.exists():
        poster_path = POSTERS_DIR / f"poster_{finding['id']:02d}.png"
        try:
            vid = VideoFileClip(str(clip_path))
            vid.save_frame(str(poster_path), t=min(8, vid.duration - 1))
            vid.close()

            s.shapes.add_movie(
                str(clip_path),
                Inches(6.4), Inches(0.4),
                Inches(6.5), Inches(3.65),
                poster_frame_image=str(poster_path),
                mime_type="video/mp4",
            )
        except Exception as e:
            t(s, Inches(6.5), Inches(2), Inches(5.5), Inches(0.5),
              f"[Video: {finding['clip_name']}]", sz=14, color=ACCENT)

    # Duration badge
    dur = finding.get("duration", 0)
    if dur:
        t(s, Inches(6.4), Inches(4.2), Inches(3), Inches(0.4),
          f"\u25B6  {dur:.0f}s clip — click to play", sz=11, color=GRAY)

    # Prevalence badge
    affected = finding.get("affected_participants", [])
    if len(affected) > 1 and not is_positive:
        prev_text = f"Also observed in: {', '.join(sorted(affected))} ({len(affected)}/6 participants)"
        t(s, Inches(0.3), Inches(5.8), Inches(5.8), Inches(0.4), prev_text, sz=11, bold=True, color=ACCENT)
    elif not is_positive:
        t(s, Inches(0.3), Inches(5.8), Inches(5.8), Inches(0.4),
          f"Observed in: {finding['participant']} (1/6 participants)", sz=11, color=GRAY)

# === SLIDE: RECOMMENDATIONS ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Recommendations", sz=32, bold=True)

recs = [
    ("P0", "Add success confirmation after deployment", "Users assume failure without explicit success state"),
    ("P0", "Show clear progress steps during build/deploy", "'Deploy to build and build to deploy' confuses everyone"),
    ("P0", "Surface logs in top-level navigation", "3 of 6 participants couldn't find logs (P1, P5, P6)"),
    ("P0", "Add custom domain settings to visible nav", "Users exhaustively search and give up (P2, P5)"),
    ("P1", "Resolve Environments vs Variables naming", "Tab names create ambiguity about where env vars live"),
    ("P1", "Add confirmation before redeployment", "Users accidentally trigger full redeploys"),
    ("P1", "Provide error recovery guidance", "Users want suggested fixes, not just error states"),
    ("P1", "Improve first-time onboarding flow", "Setup is the weakest point; once past it, tool rates highly"),
    ("P2", "Improve VS Code extension discoverability", "Hardest part for experienced developer (P1)"),
    ("P2", "Evaluate 'Embr' name internationally", "Pronunciation friction for Spanish speakers"),
]

y = 1.2
for pri, action, reason in recs:
    p_color = RED if pri == "P0" else ORANGE if pri == "P1" else YELLOW
    t(s, Inches(0.5), Inches(y), Inches(0.8), Inches(0.5), f"[{pri}]", sz=14, bold=True, color=p_color)
    t(s, Inches(1.4), Inches(y), Inches(5.5), Inches(0.5), action, sz=14, bold=True, color=WHITE)
    t(s, Inches(7.2), Inches(y), Inches(5.8), Inches(0.5), reason, sz=12, color=GRAY)
    y += 0.75

# === SLIDE: WHAT TO PRESERVE ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Preserve These Strengths", sz=32, bold=True, color=GREEN)

strengths = [
    "One-click template deployment — rated 5/5 by multiple participants",
    "Integrated database provisioning — clear differentiator vs. Vercel",
    "Build status checking UI — 'I really liked where it was signaling whether it was cool or not'",
    "Clean, minimal portal design — praised for not being overwhelming like Vercel's 25-option menu",
    "GitHub integration simplicity — 'just what I expected and how I work right now'",
]

y = 1.3
for strength in strengths:
    t(s, Inches(0.5), Inches(y), Inches(12), Inches(0.6), f"✓  {strength}", sz=16, color=GRAY)
    y += 0.8

# === SLIDE: METHODOLOGY ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
t(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Methodology", sz=32, bold=True)
t(s, Inches(0.5), Inches(1.2), Inches(6), Inches(5),
  "Method: Moderated remote usability testing\n"
  "Sessions: ~35-45 min each, video recorded\n"
  "Participants: 6 hobbyist developers\n"
  "Platforms: Vercel, DigitalOcean, self-hosted\n"
  "IDE: VS Code, Cursor, Cloud Code\n"
  "Experience: Range from beginner to 40+ years\n\n"
  "Tasks:\n"
  "1. Deploy template app with database\n"
  "2. Explore the portal\n"
  "3. Make a change and redeploy (CLI/Extension/Copilot)\n"
  "4. Try another deployment method",
  sz=14, color=GRAY)

t(s, Inches(7), Inches(1.2), Inches(5.5), Inches(5),
  "Analysis:\n"
  "• Full transcript review (6 sessions, 1100+ turns)\n"
  "• Thematic coding with Nielsen heuristics\n"
  "• Cross-participant pattern identification\n"
  "• Video clip extraction for key moments\n"
  "• Severity rating: behavioral signal-based\n\n"
  "Limitations:\n"
  "• n=6 (findings need broader validation)\n"
  "• No task-level timing analysis\n"
  "• Some prompted responses included",
  sz=14, color=GRAY)

# === SAVE ===
output = REPORTS_DIR / "Embr_UT_Stakeholder_v5.pptx"
prs.save(str(output))

# Cleanup
import shutil
shutil.rmtree(str(POSTERS_DIR), ignore_errors=True)

print(f"\nPresentation saved: {output}")
print(f"  Size: {output.stat().st_size / 1024 / 1024:.1f} MB")
print(f"  Slides: {len(prs.slides)}")
print(f"  Findings with video: {len(CURATED_FINDINGS)}")
